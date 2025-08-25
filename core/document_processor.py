import os
from typing import List
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from utils.logger import logger

class DocumentProcessor:
    
    @staticmethod
    def extract_text_from_file(file_path: str) -> str:
        """Extract text from various document formats"""
        try:
            # Clean the file path by removing any surrounding quotes
            file_path = file_path.strip('\'"')
            extension = os.path.splitext(file_path)[1].lower()
            
            if extension == '.pdf':
                return DocumentProcessor._extract_from_pdf(file_path)
            elif extension == '.docx':
                return DocumentProcessor._extract_from_docx(file_path)
            elif extension == '.pptx':
                return DocumentProcessor._extract_from_pptx(file_path)
            elif extension == '.txt':
                return DocumentProcessor._extract_from_txt(file_path)
            else:
                raise ValueError(f"Unsupported file format: {extension}")
                
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return ""
    
    @staticmethod
    def _extract_from_pdf(file_path: str) -> str:
        text = ""
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text
    
    @staticmethod
    def _extract_from_docx(file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    
    @staticmethod
    def _extract_from_pptx(file_path: str) -> str:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    @staticmethod
    def _extract_from_txt(file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()